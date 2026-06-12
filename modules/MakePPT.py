# import xlwings as wx

import io
import os
from pathlib import Path

import cv2
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from PyQt6.QtWidgets import QMessageBox

try:
    from .GetData import Futam, Horses
except:
    from GetData import Futam, Horses


class MakePPT:
    def __init__(self, mainwindow, jockeys, titles, addons=[], location=""):
        self.rome_num = [
            "I",
            "II",
            "III",
            "IV",
            "V",
            "VI",
            "VII",
            "VIII",
            "IX",
            "X",
            "XI",
            "XII",
            "XIII",
            "XIV",
        ]
        self.titles = titles
        self.jockeys = jockeys
        self.track = 11

        if location == "" and not os.path.isdir("ppt"):
            os.makedirs("ppt")
        elif not os.path.isdir(location):
            try:
                os.makedirs(location, exist_ok=True)
            except:
                location = ""

        for title in self.titles:
            ppt = Presentation()

            ppt.slide_width = Inches(10)  # Set width
            ppt.slide_height = Inches(7.5)  # Set height
            slide_layout = ppt.slide_layouts[6]

            self.slide1(ppt, slide_layout, title)
            self.slide2(ppt, slide_layout, title)

            if int(title.id) == 0:
                self.slide3(ppt, slide_layout, title)
                self.slide4(ppt, slide_layout, title, True)

            elif self.titles[-1].id == title.id:
                self.slide3(ppt, slide_layout, self.titles[int(title.id) - 1])
                self.slide4(ppt, slide_layout, self.titles[int(title.id) - 1])
                self.emptySlide(ppt, slide_layout)
                self.slide3(ppt, slide_layout, title, True)
                self.slide4(ppt, slide_layout, title, True)

            else:
                self.slide3(ppt, slide_layout, self.titles[int(title.id) - 1])
                self.slide4(ppt, slide_layout, self.titles[int(title.id) - 1])

            self.emptySlide(ppt, slide_layout)
            self.slide5(ppt, slide_layout)

            if addons:
                for addon in addons:
                    if addon.id != "":
                        self.addon_slide(ppt, slide_layout, addon)

            file_name = (
                f"{location}/{title.daily}. futam.pptx"
                if location != ""
                else f"./ppt/{title.daily}. futam.pptx"
            )

            for i, slide in enumerate(ppt.slides):
                self.set_duration(slide, 15)

            self.enable_looping(ppt)
            ppt.save(file_name)

            print(f"{file_name} created!")

        if location == "":
            QMessageBox.information(
                mainwindow,
                "PPT file creation",
                f"The powerpoint files are created to:\n{Path(__file__).resolve().parent}/ppt/",
            )
        else:
            QMessageBox.information(
                mainwindow,
                "PPT file creation",
                f"The powerpoint files are created to:\n{location}",
            )

    def addon_slide(self, ppt, slide_layout, addon):
        addonsl = ppt.slides.add_slide(slide_layout)
        # self.set_slide_duration(slide1, 5)
        # slide1.slide_show.transition.duration = 50
        addonbc = addonsl.background.fill
        addonbc.solid()
        addonbc.fore_color.rgb = RGBColor(0, 0, 0)

        if addon and addon.get_type() == "video":
            left = Inches(0)
            top = Inches(0)
            width = Inches(10)
            height = Inches(7.5)

            cap = cv2.VideoCapture(addon.path)
            success, frame = cap.read()
            cap.release()

            poster_stream = None
            if success:
                # 2. Convert the image array to a JPEG format in RAM
                is_success, buffer = cv2.imencode(".jpg", frame)
                if is_success:
                    # 3. Wrap the bytes into a file-like object
                    poster_stream = io.BytesIO(buffer)
            # Add the video
            movie = addonsl.shapes.add_movie(
                addon.path,
                left,
                top,
                width,
                height,
                poster_frame_image=poster_stream,
                mime_type=addon.type,
            )
        elif addon and addon.get_type() == "image":
            left = Inches(1)
            top = Inches(1)
            height = Inches(4.5)
            pic = addonsl.shapes.add_picture(addon.path, left, top, height=height)

    def slide1(self, ppt, slide_layout, futam):
        slide1 = ppt.slides.add_slide(slide_layout)
        # self.set_slide_duration(slide1, 5)
        # slide1.slide_show.transition.duration = 50
        slide1bc = slide1.background.fill
        slide1bc.solid()
        slide1bc.fore_color.rgb = RGBColor(0, 0, 0)

        title_box = slide1.shapes.add_textbox(
            Inches(0), Inches(0), Inches(10), Inches(1.4)
        )
        title_frame = title_box.text_frame

        title = title_frame.add_paragraph()
        # print(futam)
        title_str = f"{futam.daily}. {futam.title}".strip()

        title.text = title_str
        if len(title_str) < 39:
            title.font.size = Pt(45)
        elif len(title_str) <= 58:
            title.font.size = Pt(40)
        else:
            title.font.size = Pt(36)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)

        title.alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        text_box = slide1.shapes.add_textbox(
            Inches(4.0), Inches(1.35), Inches(6.0), Inches(1)
        )
        text_frame = text_box.text_frame

        text = text_frame.add_paragraph()
        text.text = f"Pálya: {self.track} Kincsem Park"
        text.font.size = Pt(36)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 229, 121)

        text.alignment = PP_ALIGN.RIGHT
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

        data_box = slide1.shapes.add_textbox(
            Inches(0), Inches(2.65), Inches(4.5), Inches(1.6)
        )
        data_frame = data_box.text_frame

        data_text = data_frame.add_paragraph()
        data_text.text = f"{futam.dist} m\n{futam.track}"

        data_text.font.size = Pt(36)
        data_text.font.bold = True
        data_text.font.color.rgb = RGBColor(255, 229, 121)

        data_text.alignment = PP_ALIGN.LEFT
        data_frame.word_wrap = True
        data_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

        time_box = slide1.shapes.add_textbox(
            Inches(5.72), Inches(2.65), Inches(1.78), Inches(1.60)
        )
        time_frame = time_box.text_frame

        time_text = time_frame.add_paragraph()
        time_text.text = f"Start:\n{futam.time}"

        time_text.font.size = Pt(36)
        time_text.font.bold = True
        time_text.font.color.rgb = RGBColor(255, 255, 255)

        time_text.alignment = PP_ALIGN.LEFT
        time_frame.word_wrap = True
        time_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

        slide1.shapes.add_picture(
            "clock.jpeg", Inches(7.47), Inches(3.0), Inches(1.17), Inches(1.13)
        )

        text1_box = slide1.shapes.add_textbox(
            Inches(0), Inches(4.65), Inches(4.0), Inches(1.0)
        )
        text1_frame = text1_box.text_frame

        text1 = text1_frame.add_paragraph()
        text1.text = "Véleményünk:"
        text1.font.size = Pt(36)
        text1.font.bold = True
        text1.font.color.rgb = RGBColor(255, 255, 255)

        text1.alignment = PP_ALIGN.LEFT
        text1_frame.word_wrap = True
        text1_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

        opinion_box = slide1.shapes.add_textbox(
            Inches(0), Inches(5.34), Inches(10), Inches(1.92)
        )
        opinion_frame = opinion_box.text_frame

        opinion_text = opinion_frame.add_paragraph()
        opinion_text.text = f"{futam.opinion}"

        opinion_text.font.size = Pt(30)
        opinion_text.font.bold = False
        opinion_text.font.color.rgb = RGBColor(255, 255, 255)

        time_text.alignment = PP_ALIGN.LEFT
        opinion_frame.word_wrap = True
        opinion_frame.vertical_anchor = MSO_ANCHOR.TOP

    def slide2(self, ppt, slide_layout, futam):
        jockeys_list = [jockey for jockey in self.jockeys if jockey.Fnum == futam.id]

        slide2 = ppt.slides.add_slide(slide_layout)
        # slide2.slide_show.transition.duration = 50
        slide2bc = slide2.background.fill
        slide2bc.solid()
        slide2bc.fore_color.rgb = RGBColor(0, 0, 0)

        horse_box = slide2.shapes.add_textbox(
            Inches(0), Inches(0), Inches(5.25), Inches(7.5)
        )
        horse_frame = horse_box.text_frame

        jockey_box = slide2.shapes.add_textbox(
            Inches(5.25), Inches(0), Inches(5.25), Inches(7.5)
        )
        jockey_frame = jockey_box.text_frame

        for row in jockeys_list:
            # horse side

            horse = horse_frame.add_paragraph()
            horse.text = f"{row.Hnum}. {row.Hname.upper()}"
            # print(f"{row.horse_number}. {row.horse_name.upper()}")
            horse.font.size = Pt(30)
            horse.font.bold = True
            horse.font.color.rgb = RGBColor(255, 229, 121)

            horse.alignment = PP_ALIGN.LEFT
            horse_frame.word_wrap = True
            horse_frame.vertical_anchor = MSO_ANCHOR.TOP

            # Jockey side

            jockey = jockey_frame.add_paragraph()
            jockey.text = f"{row.DJname}"
            jockey.font.size = Pt(30)
            jockey.font.bold = False
            jockey.font.color.rgb = RGBColor(255, 255, 255)

            jockey.alignment = PP_ALIGN.LEFT
            jockey_frame.word_wrap = True
            jockey_frame.vertical_anchor = MSO_ANCHOR.TOP

    def slide3(self, ppt, slide_layout, futam, hide=False):
        slide3 = ppt.slides.add_slide(slide_layout)
        # slide3.slide_show.transition.duration = 50
        slide3bc = slide3.background.fill
        slide3bc.solid()
        slide3bc.fore_color.rgb = RGBColor(0, 0, 0)

        title_box = slide3.shapes.add_textbox(
            Inches(0), Inches(0), Inches(10), Inches(1.4)
        )
        title_frame = title_box.text_frame

        title = title_frame.add_paragraph()

        title_str = f"{futam.daily}. {futam.title}".strip()

        title.text = title_str
        if len(title_str) < 39:
            title.font.size = Pt(45)
        elif len(title_str) <= 58:
            title.font.size = Pt(40)
        else:
            title.font.size = Pt(36)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)

        title.alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        text_box = slide3.shapes.add_textbox(
            Inches(0), Inches(1.4), Inches(8.25), Inches(2)
        )
        text_frame = text_box.text_frame

        text = text_frame.add_paragraph()
        text.text = f"Pálya: {self.track} Kincsem Park\nBefutási sorrend:"
        text.font.size = Pt(48)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 229, 121)

        text.alignment = PP_ALIGN.LEFT
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        order_box = slide3.shapes.add_textbox(
            Inches(0), Inches(3.37), Inches(10), Inches(2.52)
        )
        order_frame = order_box.text_frame

        order = order_frame.add_paragraph()
        order.text = "I.\nII.\nIII."
        order.font.size = Pt(48)
        order.font.bold = True
        order.font.color.rgb = RGBColor(255, 255, 255)

        order.alignment = PP_ALIGN.LEFT
        order_frame.word_wrap = True
        order_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        time_box = slide3.shapes.add_textbox(
            Inches(0), Inches(5.9), Inches(3.39), Inches(0.91)
        )
        time_frame = time_box.text_frame

        time_text = time_frame.add_paragraph()
        time_text.text = "Idő:"

        time_text.font.size = Pt(48)
        time_text.font.bold = True
        time_text.font.color.rgb = RGBColor(255, 229, 121)

        time_text.alignment = PP_ALIGN.LEFT
        time_frame.word_wrap = True
        time_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

    def slide4(self, ppt, slide_layout, futam, hide=False):
        slide4 = ppt.slides.add_slide(slide_layout)
        # slide4.slide_show.transition.duration = 50
        slide4bc = slide4.background.fill
        slide4bc.solid()
        slide4bc.fore_color.rgb = RGBColor(0, 0, 0)

        title_box = slide4.shapes.add_textbox(
            Inches(0), Inches(0), Inches(10), Inches(1.4)
        )
        title_frame = title_box.text_frame
        title = title_frame.add_paragraph()

        title_str = f"{futam.daily}. {futam.title}".strip()

        title.text = title_str
        if len(title_str) < 39:
            title.font.size = Pt(45)
        elif len(title_str) <= 58:
            title.font.size = Pt(40)
        else:
            title.font.size = Pt(36)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)

        title.alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        text_box = slide4.shapes.add_textbox(
            Inches(0), Inches(1.5), Inches(10), Inches(1.72)
        )
        text_frame = text_box.text_frame

        text = text_frame.add_paragraph()
        text.text = f"Pálya: {self.track} Kincsem Park\nBefutási sorrend:"
        text.font.size = Pt(48)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 229, 121)

        text.alignment = PP_ALIGN.LEFT
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        dividend_box = slide4.shapes.add_textbox(
            Inches(0), Inches(3.57), Inches(10), Inches(2.52)
        )
        dividend_frame = dividend_box.text_frame

        dividend = dividend_frame.add_paragraph()
        dividend.text = f"Tét:{'\t' * 6}1\nHely:{'\t' * 6}1\nBefutó:{'\t' * 4}1\nHbefutó:{'\t' * 4}1"
        dividend.font.size = Pt(48)
        dividend.font.bold = True
        dividend.font.color.rgb = RGBColor(255, 255, 255)

        dividend.alignment = PP_ALIGN.LEFT
        dividend_frame.word_wrap = True
        dividend_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def slide5(self, ppt, slide_layout):
        slide5 = ppt.slides.add_slide(slide_layout)
        # slide5.slide_show.transition.duration = 50
        slide5bc = slide5.background.fill
        slide5bc.solid()
        slide5bc.fore_color.rgb = RGBColor(0, 0, 0)

        title_box = slide5.shapes.add_textbox(
            Inches(0), Inches(0), Inches(10), Inches(2.3)
        )
        title_frame = title_box.text_frame

        title = title_frame.add_paragraph()
        title.text = "Nem hivatalos\n  befutási sorrend:"
        title.font.size = Pt(78)
        title.font.bold = True
        title.font.color.rgb = RGBColor(255, 229, 121)

        title.alignment = PP_ALIGN.LEFT
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        text_box = slide5.shapes.add_textbox(
            Inches(0), Inches(3.75), Inches(10), Inches(2.21)
        )
        text_frame = text_box.text_frame

        text = text_frame.add_paragraph()
        text.text = " –  – "
        text.font.size = Pt(125)
        text.font.bold = True
        text.font.color.rgb = RGBColor(255, 255, 255)

        text.alignment = PP_ALIGN.CENTER
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def emptySlide(self, ppt, slide_layout):
        empty_slide = ppt.slides.add_slide(slide_layout)
        # slide5.slide_show.transition.duration = 50
        empty_slidebc = empty_slide.background.fill
        empty_slidebc.solid()
        empty_slidebc.fore_color.rgb = RGBColor(0, 0, 0)

    def set_duration(self, slide, seconds):
        """Sets the 'Advance After' time with strict LibreOffice compatibility."""
        ms = int(seconds * 1000)

        # Check if transition already exists
        transition = slide.element.find(qn("p:transition"))

        if transition is None:
            # We MUST include a transition type (like <p:fade/>) for LibreOffice
            # We also set advClick="0" to prioritize the timer over mouse clicks
            new_xml = (
                f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                f'spd="fast" advClick="0" advTm="{ms}">'
                f"<p:cut />"
                f"</p:transition>"
            )
            transition = parse_xml(new_xml)
            # The transition tag MUST be appended at the end of the slide element
            slide.element.append(transition)
        else:
            transition.set("advTm", str(ms))
            transition.set("advClick", "0")
            # If the existing transition is empty, add a fade effect
            if len(transition.getchildren()) == 0:
                transition.append(
                    parse_xml(
                        '<p:cut xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
                    )
                )

    def enable_looping(self, pres):
        # 1. Access the Presentation Properties part
        try:
            prs_props_part = pres.part.part_related_by(RT.PRES_PROPS)
        except KeyError:
            # If the part doesn't exist, this script won't work easily
            # (usually it exists in any standard saved PPTX)
            return

        # 2. Get the root XML element (p:presentationPr)
        presentationPr = parse_xml(prs_props_part.blob)

        # 3. Find or create the <p:showPr> element and set loop="1"
        # The namespace for 'p' is usually required for the search
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        show_pr = presentationPr.find(".//p:showPr", namespaces=nsmap)

        if show_pr is None:
            # Create it if it's missing (rare, but possible)
            from lxml import etree

            show_pr = etree.SubElement(presentationPr, "{%s}showPr" % nsmap["p"])

        show_pr.set("loop", "1")

        # 4. Write the modified XML back into the part
        prs_props_part._blob = etree.tostring(presentationPr)


if __name__ == "__main__":
    titles = []
    jockeys = []
    with open("./csv/titles_data.csv", "r", encoding="utf-8") as f:
        fs = f.readline()
        for ln in f:
            # print(ln.strip())
            titles.append(Futam(ln))

    with open("./csv/jockeys_data.csv", "r", encoding="utf-8") as f:
        fs = f.readline()
        for ln in f:
            jockeys.append(Horses(ln))

    MakePPT(jockeys, titles)
    print("Id;Daily;Title;Distance;Start time;Start type;Opinion")
    for i in titles:
        print(i)
